from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Crear presentación
prs = Presentation()

# Función para agregar diapositiva con título y puntos
def agregar_diapositiva(titulo, puntos):
    slide_layout = prs.slide_layouts[1]  # Título y contenido
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = titulo
    content.text = "\n".join(puntos)

# Diapositiva 1: Introducción
agregar_diapositiva("¿Qué es la Visión por Computadora?", [
    "Rama de la Inteligencia Artificial.",
    "Permite que las computadoras 'vean' e interpreten imágenes o video.",
    "Se usa para reconocer objetos, seguir movimiento, identificar rostros, etc."
])

# Diapositiva 2: ¿Cómo ve una computadora?
agregar_diapositiva("¿Cómo ve una computadora?", [
    "Las imágenes son matrices de píxeles (números).",
    "Escala de grises: matriz 2D con valores entre 0-255.",
    "Color: matriz 3D con canales BGR (Blue, Green, Red)."
])

# Diapositiva 3: ¿Qué es OpenCV?
agregar_diapositiva("OpenCV y la librería cv2", [
    "Open Source Computer Vision Library.",
    "Usada con Python como 'cv2'.",
    "Permite leer, procesar, analizar y guardar imágenes y video."
])

# Diapositiva 4: Funciones básicas en cv2
agregar_diapositiva("Operaciones básicas con cv2", [
    "Leer y mostrar imágenes.",
    "Convertir a escala de grises.",
    "Guardar imágenes procesadas.",
    "Capturar desde cámara.",
    "Redimensionar, recortar, modificar color."
])

# Diapositiva 5: Ejemplo práctico
agregar_diapositiva("Ejemplo práctico: flujo de procesamiento", [
    "1. Cargar imagen",
    "2. Convertir a escala de grises",
    "3. Detectar bordes (Canny)",
    "4. Dibujar contornos",
    "5. Guardar resultados"
])

# Diapositiva 6: Código Python con OpenCV
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Código de ejemplo (cv2)"

code = '''import cv2
img = cv2.imread("urban.jpg")
gris = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
bordes = cv2.Canny(gris, 100, 200)
contornos, _ = cv2.findContours(bordes, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
cv2.drawContours(img, contornos, -1, (0,255,0), 2)
cv2.imshow("Contornos", img)
cv2.waitKey(0)
cv2.destroyAllWindows()'''

content = slide.placeholders[1]
content.text = code

# Guardar archivo
prs.save("Vision_Por_Computadora.pptx")
print("Presentación creada: Vision_Por_Computadora.pptx")
